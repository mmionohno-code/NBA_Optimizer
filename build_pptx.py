"""
Build PowerPoint presentation for the NBA Optimizer project.
Tells the full story: Problem → Data → Method → Results → Limitations.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import os

print("Building PowerPoint presentation...")

# ── Colors ──────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x2C, 0x2C, 0x2C)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
BLUE = RGBColor(0x2C, 0x80, 0xD1)

# ── Load data for stats ─────────────────────────────────────────────────────
df_all = pd.read_csv('nba_clustered.csv')
df_2324 = df_all[df_all['SEASON'] == '2023-24'].copy()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ─────────────────────────────────────────────────────────

def add_background(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_slide(slide, left, top, width, height, bullets, font_size=16,
                     color=WHITE, bold_first=False, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        if bold_first and ':' in bullet:
            # Can't do partial bold easily in python-pptx simple mode
            p.font.bold = False
    return tf

def add_image_safe(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                  Inches(width), Inches(height))
        return True
    return False

# ============================================================
# SLIDE 1 — TITLE
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide, NAVY)

add_textbox(slide, 1, 1.8, 11.3, 1.5,
            "NBA ROSTER OPTIMIZATION ENGINE",
            font_size=40, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 3.3, 11.3, 0.8,
            "Building the Statistically Optimal Team Under a Salary Cap",
            font_size=22, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 4.5, 11.3, 0.6,
            "Data-Driven  |  PCA Composite Scoring  |  K-Means Archetypes  |  MILP Optimization  |  Synergy Model",
            font_size=14, color=GOLD, bold=False, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.0, 11.3, 0.5,
            f"1,185 Player-Seasons  |  9 Validated Stats  |  6 Archetypes  |  10 Optimized Rosters",
            font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

print("  Slide 1 (Title) done")

# ============================================================
# SLIDE 2 — THE PROBLEM
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "THE PROBLEM", font_size=32, color=GOLD, bold=True)

bullets = [
    "NBA teams spend $150M+ on rosters — but how do you know you're spending it right?",
    "Traditional roster building relies on scouts, gut feelings, and media narratives",
    "Players are evaluated in isolation — but basketball is a team sport where fit matters",
    "No single stat captures a player's total value — points ignore defense, rebounds ignore playmaking",
    "",
    "GOAL: Build a mathematically optimal 15-man NBA roster that maximizes total player value",
    "under a real salary cap constraint, using only publicly available data and zero human bias.",
]
add_bullet_slide(slide, 0.8, 1.4, 11.7, 5.5, bullets, font_size=18, color=WHITE)

print("  Slide 2 (Problem) done")

# ============================================================
# SLIDE 3 — DATA OVERVIEW
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "DATA COLLECTION", font_size=32, color=GOLD, bold=True)

bullets = [
    "Source: NBA Stats API (nba_api) — official NBA.com backend. Same data NBA teams see.",
    "3 seasons: 2021-22, 2022-23, 2023-24 (~400 players per season = 1,185 player-seasons)",
    "Quality filter: minimum 20 games played and 10 minutes per game",
    "Recency weights: 2021-22 (20%), 2022-23 (35%), 2023-24 (45%)",
    "Salary data: 2023-24 verified contracts + 87 manual corrections for imputation errors",
    "ON/OFF splits from 5-man lineup data (LeagueDashLineups API)",
    "2-man lineup data for pairwise synergy computation (6,000+ pair-seasons)",
]
add_bullet_slide(slide, 0.8, 1.4, 11.7, 5.5, bullets, font_size=17, color=WHITE)

print("  Slide 3 (Data) done")

# ============================================================
# SLIDE 4 — FEATURE ENGINEERING
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "FEATURE ENGINEERING — Cleaning the Signal", font_size=30, color=GOLD, bold=True)

bullets = [
    "Team Context Adjustment: subtract team average from each stat to isolate individual contribution",
    "Bayesian TS% Shrinkage: pull low-volume shooters toward league average (prevents small-sample inflation)",
    "FG3% Neutralization: non-shooters (<0.5 FG3A/game) set to league avg — 'no info' not 'bad shooter'",
    "W_PCT Residualization: OLS regression removes team-quality bias from ON_OFF_DIFF",
    "INFLUENCE_SCORE = USG% x TS% — separates efficient volume scorers from inefficient ones",
    "IS_SHOOTER flag: created before neutralization so optimizer knows who actually shoots threes",
]
add_bullet_slide(slide, 0.8, 1.4, 11.7, 5.5, bullets, font_size=16, color=WHITE)

print("  Slide 4 (Feature Engineering) done")

# ============================================================
# SLIDE 5 — CORRELATION ANALYSIS
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "STAT SELECTION — Three-Framework Correlation Analysis", font_size=28, color=GOLD, bold=True)

bullets = [
    "Framework 1: Team-adjusted stats vs Win% — which stats predict team winning?",
    "Framework 2: Individual counting stats vs Win% — which individual skills tie to winning?",
    "Framework 3: All stats vs Plus/Minus (team-adjusted) — captures invisible contributions",
    "Spearman Validation: re-test everything with rank correlation to catch non-linear relationships",
    "Selection Rule: stat included if p < 0.05 in at least one framework AND validated by Spearman",
    "DREB_PCT eliminated: position-group adjustment revealed it was measuring height, not skill",
    "",
    "RESULT: 9 stats survived — OFF_RATING_ADJ, DEF_RATING_ADJ, AST_PCT_ADJ, BLK, STL, TS%, FG3%, INFLUENCE_SCORE, ON_OFF_DIFF",
]
add_bullet_slide(slide, 0.8, 1.4, 11.7, 5.5, bullets, font_size=15, color=WHITE)

print("  Slide 5 (Correlation) done")

# ============================================================
# SLIDE 6 — INTERCORRELATION CHART
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 5.5, 0.8,
            "INTERCORRELATION CHECK", font_size=28, color=GOLD, bold=True)

add_textbox(slide, 0.8, 1.2, 5.0, 1.2,
            "Before entering PCA, all 9 stats must be sufficiently independent. "
            "No pair exceeds |r| > 0.7 — each stat captures a unique dimension of player skill. "
            "PIE_ADJUSTED was removed due to r=0.78 with INFLUENCE_SCORE (redundant).",
            font_size=14, color=LIGHT_GRAY)

add_image_safe(slide, 'charts/intercorrelation_final.png', 6.2, 0.4, 6.5, 6.5)

print("  Slide 6 (Intercorrelation) done")

# ============================================================
# SLIDE 7 — PCA
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "PCA — Deriving the Composite Score", font_size=30, color=GOLD, bold=True)

add_textbox(slide, 0.8, 1.3, 5.5, 3.0,
            "PCA (Principal Component Analysis) takes 9 stats per player and compresses them "
            "into a single score (0-100).\n\n"
            "The weights are not chosen by us — they are derived entirely from the data.\n\n"
            "PC1 alone captures ~30% of variance. We use PC1+PC2+PC3 (65% combined), "
            "weighted by their variance-explained shares:\n"
            "  PC1: 47%  |  PC2: 29%  |  PC3: 24%\n\n"
            "This ensures the dominant basketball excellence signal (PC1) carries the most "
            "weight while PC2/PC3 add independent information.",
            font_size=14, color=WHITE)

add_image_safe(slide, 'charts/pca_weights_final.png', 6.5, 1.0, 6.3, 3.0)
add_image_safe(slide, 'charts/pca_variance_explained.png', 6.5, 4.1, 6.3, 3.0)

print("  Slide 7 (PCA) done")

# ============================================================
# SLIDE 8 — TOP 20 RANKINGS
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 5.5, 0.8,
            "TOP 20 PLAYERS — 2023-24", font_size=28, color=GOLD, bold=True)

top20 = df_2324.sort_values('COMPOSITE_SCORE_NORM', ascending=False).head(20)
rank_text = ""
for i, (_, row) in enumerate(top20.iterrows()):
    name = row['PLAYER_NAME']
    score = row['COMPOSITE_SCORE_NORM']
    rank_text += f"{i+1:>2}. {name:<28} {score:.1f}\n"

add_textbox(slide, 0.8, 1.3, 5.0, 5.5, rank_text,
            font_size=13, color=WHITE, font_name='Consolas')

add_image_safe(slide, 'charts/top20_players_final.png', 6.3, 0.8, 6.5, 6.0)

print("  Slide 8 (Top 20) done")

# ============================================================
# SLIDE 9 — K-MEANS CLUSTERING
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "K-MEANS CLUSTERING — Discovering Player Archetypes", font_size=28, color=GOLD, bold=True)

add_textbox(slide, 0.8, 1.2, 5.5, 2.5,
            "K-Means groups 1,185 player-seasons into clusters based on their 9-stat profiles. "
            "Players with similar statistical fingerprints land in the same group.\n\n"
            "Elbow method validated K=7 as optimal (7 clusters mapped to 6 archetype labels).\n\n"
            "Dynamic labeling: archetypes assigned by centroid statistics, not hardcoded cluster numbers. "
            "This prevents labels from breaking when data changes.",
            font_size=14, color=WHITE)

add_image_safe(slide, 'charts/elbow_method.png', 6.3, 0.4, 6.5, 3.3)
add_image_safe(slide, 'charts/cluster_visualization.png', 6.3, 3.8, 6.5, 3.5)

arch_text = (
    "6 ARCHETYPES:\n"
    "  Elite Playmaker    — Jokic, SGA, Luka, Embiid, Giannis\n"
    "  Two-Way Big        — Wembanyama, AD, Gobert\n"
    "  Versatile Scorer   — Franz Wagner, Garland, Randle\n"
    "  Defensive Wing     — OG Anunoby, KCP, Aaron Gordon\n"
    "  Perimeter Scorer   — Markkanen, RJ Barrett\n"
    "  Bench / Role Player — Depth pieces"
)
add_textbox(slide, 0.8, 3.8, 5.5, 3.0, arch_text,
            font_size=13, color=LIGHT_GRAY, font_name='Consolas')

print("  Slide 9 (Clustering) done")

# ============================================================
# SLIDE 10 — RADAR CHARTS
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.6,
            "ARCHETYPE STAT PROFILES", font_size=28, color=GOLD, bold=True)

add_textbox(slide, 0.8, 1.0, 11.7, 0.5,
            "Each archetype has a distinct statistical fingerprint — radar charts show average stat profile on a 0-1 normalized scale.",
            font_size=14, color=LIGHT_GRAY)

add_image_safe(slide, 'charts/archetype_radar_charts.png', 0.5, 1.6, 12.3, 5.5)

print("  Slide 10 (Radar Charts) done")

# ============================================================
# SLIDE 11 — OPTIMIZER (MILP)
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "MILP OPTIMIZATION — Building the Roster", font_size=30, color=GOLD, bold=True)

bullets = [
    "MILP (Mixed Integer Linear Programming): mathematically proven optimal solution, not heuristic",
    "399 binary decision variables: each player is either on the roster (1) or not (0)",
    "Objective: maximize total composite score (or variant objectives per scenario)",
    "",
    "HARD CONSTRAINTS:",
    "  - Salary cap (varies by scenario: $90M to $165M)",
    "  - Exactly 15 players on the roster",
    "  - Archetype minimums (at least 1 Elite Playmaker, 2 Two-Way Bigs, etc.)",
    "  - Max 2 players from any single NBA team (team diversity)",
    "  - Max 5 big men on roster",
    "",
    "VARIANT CONSTRAINTS (per scenario):",
    "  - Rebuild: at least 9 players age 24 or younger",
    "  - Win-Now: at least 9 veterans age 28 or older",
    "  - Three-Point Era: at least 10 real shooters (IS_SHOOTER flag)",
    "  - Small Ball: max 1 Two-Way Big",
]
add_bullet_slide(slide, 0.8, 1.2, 11.7, 6.0, bullets, font_size=14, color=WHITE, spacing=Pt(3))

print("  Slide 11 (Optimizer) done")

# ============================================================
# SLIDE 12 — SYNERGY MODEL
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "SYNERGY MODEL — Player Chemistry", font_size=30, color=GOLD, bold=True)

bullets = [
    "PROBLEM: The base optimizer treats players independently. Basketball is a team sport.",
    "Some pairs elevate each other (Jokic + shooters). Some pairs clash.",
    "",
    "THREE SYNERGY LAYERS:",
    "",
    "Layer 1 — Player Profile: each player's average NET_SYNERGY across all observed 2-man pairs.",
    "    Players who consistently elevate teammates get a score boost. Weight = 0.4.",
    "",
    "Layer 2 — Pairwise Bonus: if two specific players with observed positive synergy are both",
    "    selected, the optimizer gets a bonus. Linearized via binary y_ij variables in MILP.",
    "",
    "Layer 3 — Defensive Coverage Zones: hard constraints requiring paint protection (BLK >= 1.5),",
    "    perimeter disruption (STL >= 1.3), and switchability (4+ defensive-archetype players).",
    "",
    "VALIDATION: W_NET_SYNERGY has highest correlation with Win% (r=+0.568) across all 3 seasons.",
    "The data chose NET_SYNERGY — not us.",
]
add_bullet_slide(slide, 0.8, 1.2, 11.7, 6.0, bullets, font_size=14, color=WHITE, spacing=Pt(3))

print("  Slide 12 (Synergy) done")

# ============================================================
# SLIDE 13 — 10 SCENARIOS OVERVIEW
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "10 OPTIMIZATION SCENARIOS", font_size=30, color=GOLD, bold=True)

scenario_text = (
    "Scenario A — Hard Cap ($136M)          Standard baseline roster\n"
    "Scenario B — Luxury Tax ($165M)        Best roster money can buy\n"
    "Scenario C — Budget Team ($90M)        Small-market constraint\n"
    "Scenario D — Rebuild Mode ($136M)      Min 9 players age <= 24\n"
    "Scenario E — Win-Now ($165M)           Min 9 veterans age >= 28\n"
    "Scenario F — Defensive Identity        Maximize defensive rating\n"
    "Scenario G — Offensive Identity        Maximize offensive rating\n"
    "Scenario H — Three-Point Era           Min 10 real 3PT shooters\n"
    "Scenario I — Small Ball                Max 1 Two-Way Big\n"
    "Scenario J — Value / Efficiency        Maximize VORPD (Moneyball)\n"
)
add_textbox(slide, 0.8, 1.3, 11.7, 3.0, scenario_text,
            font_size=16, color=WHITE, font_name='Consolas')

add_textbox(slide, 0.8, 5.0, 11.7, 1.0,
            "Each scenario solves a different real-world GM question.\n"
            "All 10 produce mathematically optimal rosters — proven by the MILP solver, not approximated.",
            font_size=15, color=LIGHT_GRAY)

print("  Slide 13 (Scenarios) done")

# ============================================================
# SLIDE 14 — SAMPLE ROSTER (Scenario A)
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "SCENARIO A — Hard Cap Roster ($136M)", font_size=28, color=GOLD, bold=True)

try:
    roster_a = pd.read_csv('optimized_roster_syn_A.csv')
    score_col = 'COMPOSITE_SYNERGY' if 'COMPOSITE_SYNERGY' in roster_a.columns else 'COMPOSITE_SCORE_NORM'
    roster_a = roster_a.sort_values(score_col, ascending=False)

    header = f"{'#':>2}  {'Player':<28} {'Archetype':<22} {'Score':>6}  {'Salary':>12}\n"
    header += "-" * 78 + "\n"
    for i, (_, row) in enumerate(roster_a.iterrows()):
        header += (f"{i+1:>2}  {row['PLAYER_NAME']:<28} {row['ARCHETYPE']:<22} "
                   f"{row[score_col]:>6.1f}  ${row['SALARY']:>11,.0f}\n")

    total_sal = roster_a['SALARY'].sum()
    total_score = roster_a[score_col].sum()
    header += f"\n{'TOTAL':<33} {'':22} {total_score:>6.1f}  ${total_sal:>11,.0f}"

    add_textbox(slide, 0.5, 1.2, 12.3, 6.0, header,
                font_size=12, color=WHITE, font_name='Consolas')
except Exception as e:
    add_textbox(slide, 0.8, 2.0, 11.7, 1.0, f"Roster file not available: {e}",
                font_size=16, color=RED)

print("  Slide 14 (Roster A) done")

# ============================================================
# SLIDE 15 — BUGS FOUND & FIXED
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "BUGS FOUND & FIXED — Quality Assurance", font_size=28, color=GOLD, bold=True)

bullets = [
    "87 salary imputation errors: Jimmy Butler showed $1.1M instead of $45.2M. Manually verified all 87.",
    "74 duplicate player-season rows: NBA API artifact. Same player selected twice on roster. Fixed with dedup.",
    "Hardcoded cluster labels broke on every data change: Jokic labeled 'Bench Player'. Fixed with dynamic centroid labeling.",
    "Non-shooters passed 3PT constraint: FG3% neutralization set Gobert to league avg, making him a 'shooter'. Fixed with IS_SHOOTER flag.",
    "Stale archetype name 'Two-Way Guard': optimizer constraint silently matched 0 players. Renamed to 'Versatile Scorer'.",
    "Equal PCA weights: LaMelo ranked #4 over Embiid. Fixed with variance-explained weighting.",
    "Stale methodology text: dashboard described old approach. Corrected.",
    "",
    "AUTOMATED VERIFICATION: 175 checks across all pipeline outputs — all passing.",
]
add_bullet_slide(slide, 0.8, 1.2, 11.7, 5.8, bullets, font_size=14, color=WHITE, spacing=Pt(4))

print("  Slide 15 (Bugs) done")

# ============================================================
# SLIDE 16 — LIMITATIONS (HONEST)
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "LIMITATIONS — What the Model Cannot Do", font_size=28, color=GOLD, bold=True)

bullets = [
    "1. Salary Approximation: older-season salaries estimated via cap-ratio scaling, not actual contracts.",
    "",
    "2. Defensive Measurement: DEF_RATING is shared across all 5 players on court — can't assign individual credit.",
    "   PCA weight of 0.039 reflects this honest limitation. STL and BLK carry the defensive signal instead.",
    "",
    "3. Advanced Defensive Stats: D-EPM and D-LEBRON require proprietary tracking data (Second Spectrum).",
    "   Public APIs don't expose this. NBA front offices have better defensive data than we do.",
    "",
    "4. Team Context: residualization reduces but can't fully eliminate teammate-quality contamination.",
    "   Role players on elite teams still get a slight bump.",
    "",
    "5. No Future Projection: model scores past performance, not future potential. Can't predict breakouts,",
    "   aging curves, or injuries. Builds the best team from what already happened.",
]
add_bullet_slide(slide, 0.8, 1.2, 11.7, 5.8, bullets, font_size=14, color=WHITE, spacing=Pt(2))

print("  Slide 16 (Limitations) done")

# ============================================================
# SLIDE 17 — PIPELINE ARCHITECTURE
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "PIPELINE ARCHITECTURE — 9 Scripts, End to End", font_size=28, color=GOLD, bold=True)

pipeline_text = (
    "fix_salaries.py          Correct 87 salary imputation errors in master data\n"
    "       |\n"
    "01_feature_engineering   Team adjustment, Bayesian shrinkage, correlation analysis,\n"
    "       |                 PCA scoring, VORPD computation\n"
    "       v\n"
    "02_clustering            K-Means (K=7) on 9 stats -> 7 clusters\n"
    "       |\n"
    "03_fix_labels            Dynamic centroid-based archetype labeling -> 6 names\n"
    "       |\n"
    "04_optimizer_expanded    Pre-synergy MILP optimizer (10 scenarios)\n"
    "       |\n"
    "05_compute_synergy       2-man lineup synergy (DEF + OFF + NET)\n"
    "       |\n"
    "06_validate_synergy      Statistical validation (W_NET_SYNERGY r=+0.568 vs W_PCT)\n"
    "       |\n"
    "07_optimizer_synergy     Final MILP with 3-layer synergy model (10 scenarios)\n"
    "       |\n"
    "08_build_dashboard       Excel dashboard (8 tabs) + Charts\n"
    "       |\n"
    "09_verify_pipeline       175 automated checks — all passing\n"
)
add_textbox(slide, 1.5, 1.2, 10.3, 5.8, pipeline_text,
            font_size=14, color=WHITE, font_name='Consolas')

print("  Slide 17 (Pipeline) done")

# ============================================================
# SLIDE 18 — TOOLS & TECHNOLOGIES
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.8, 0.4, 11.7, 0.8,
            "TOOLS & TECHNOLOGIES", font_size=30, color=GOLD, bold=True)

bullets = [
    "Python 3 — primary language for all pipeline scripts",
    "pandas / NumPy — data manipulation and numerical computation",
    "scikit-learn — PCA, K-Means clustering, StandardScaler",
    "SciPy — Pearson and Spearman correlation analysis",
    "PuLP (CBC solver) — Mixed Integer Linear Programming optimization",
    "matplotlib / seaborn — statistical visualization (7 charts)",
    "openpyxl — Excel dashboard generation (8 tabs)",
    "nba_api — official NBA Stats API data extraction",
    "python-pptx — this presentation",
]
add_bullet_slide(slide, 0.8, 1.4, 11.7, 5.5, bullets, font_size=18, color=WHITE)

print("  Slide 18 (Tools) done")

# ============================================================
# SLIDE 19 — CLOSING
# ============================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 1, 2.2, 11.3, 1.0,
            "NBA ROSTER OPTIMIZATION ENGINE",
            font_size=36, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 3.5, 11.3, 0.8,
            "1,185 players scored  |  175 automated checks passed  |  10 optimized rosters",
            font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 4.8, 11.3, 0.6,
            "Public data. Zero bias. Mathematically proven optimal.",
            font_size=20, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

print("  Slide 19 (Closing) done")

# ============================================================
# SAVE
# ============================================================

output_file = 'NBA_Optimizer_Presentation.pptx'
prs.save(output_file)

print(f"\n{'=' * 60}")
print(f"PRESENTATION SAVED: {output_file}")
print(f"{'=' * 60}")
print(f"Total slides: 19")
print("  1. Title")
print("  2. The Problem")
print("  3. Data Collection")
print("  4. Feature Engineering")
print("  5. Stat Selection (Correlation)")
print("  6. Intercorrelation Check + Chart")
print("  7. PCA Composite Score + Charts")
print("  8. Top 20 Rankings + Chart")
print("  9. K-Means Clustering + Charts")
print(" 10. Archetype Radar Charts")
print(" 11. MILP Optimization")
print(" 12. Synergy Model")
print(" 13. 10 Scenarios Overview")
print(" 14. Sample Roster (Scenario A)")
print(" 15. Bugs Found & Fixed")
print(" 16. Limitations")
print(" 17. Pipeline Architecture")
print(" 18. Tools & Technologies")
print(" 19. Closing")
